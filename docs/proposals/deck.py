#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""earthus 협업제안서 → 발표용 PPTX (16:9)

쓰는 법:  python3 docs/proposals/deck.py

⚠️ 이 기계에는 pptxgenjs·LibreOffice·pdftoppm 이 없다.
   그래서 python-pptx 로 만들고, **같은 자료로 HTML 미리보기를 함께 만들어**
   크롬으로 찍어 눈으로 검수한다 (docs/talk 에서 쓰던 방식과 같다).
   미리보기와 PPTX 는 같은 SLIDES 목록에서 나오므로 서로 어긋나지 않는다.

⚠️ 색은 제품에서 가져왔다. app.css 의 --teal(#3fc7c0) 과 태풍 주황,
   그리고 지구본 배경의 남색이다. 화면캡처가 전부 어두워서
   밝은 바탕에 얹으면 그림만 검은 사각형으로 떠 보인다.
"""
import io, os, re, subprocess

HERE = os.path.dirname(os.path.abspath(__file__))
IMG  = os.path.join(HERE, 'img')
PPTX = os.path.join(HERE, 'earthus-협업제안서-2026.pptx')
PREV = os.path.join(HERE, 'deck-preview.html')

# ── 색 ─────────────────────────────────────────────────────────
INK   = '0A1220'   # 우주 남색 — 바탕 (지배색)
INK2  = '111C2E'   # 카드 바탕
TEAL  = '3FC7C0'   # 제품 강조색
AMBER = 'FF9F43'   # 태풍 주황 — 경고·강조
PAPER = 'F2F5F8'
MUTE  = '8A97A8'
LINE  = '25334A'

W, H = 13.333, 7.5
M    = 0.9

# 슬라이드 자료 —  kind 별로 그리는 법이 다르다
SLIDES = [
 dict(kind='cover', n='', title='earthus 협업제안서',
      sub='서울창조경제혁신센터 오픈 이노베이션 2026',
      sub2='SK텔레콤 · LS ELECTRIC · 두산에너빌리티   |   달루어 · 2026년 8월',
      img='globe-cyclone.png'),

 dict(kind='thesis', n='', title='예보는 발표되고 사라집니다.',
      body='earthus 는 한·일·영·미 기상청 예보를 매시간 박제해,\n'
           '무엇이 맞고 틀렸는지를 3D 지구본 위에서 보여줍니다.',
      foot='2026년 8월 4일 정식 오픈 · earthus.net 운영 중'),

 dict(kind='stats', n='1', title='가진 것',
      note='숫자는 전부 2026년 8월 6일 실측입니다',
      stats=[('54', '개', '자료 수집 파이프라인\nAWS Lambda · 자동 실행'),
             ('18', '개', '자료 출처 기관\n4개국 기상청 포함'),
             ('58', '개', '지구본 레이어\n7가지 표현 방식'),
             ('10.2', 'MB', '하루 축적량\n11일 연속 편차 3% 이내')],
      body='출처 18곳 — 기상청(ASOS·AWS·동네예보·특보·낙뢰·해양·고층), 일본 기상청(AMeDAS·경보),\n'
           'NOAA(GMGSI 구름합성 · NDBC/OSMC 부이 · SWPC 우주기상 · METAR · NWS 쓰나미),\n'
           'NASA(GIBS 위성타일 · FIRMS 화재), ECMWF, Copernicus CAMS, USGS 지진, GDACS 재난,\n'
           'GDELT 사건, Celestrak 위성궤도, Open-Meteo, 해양수산부(KHOA), 산림청, 영국 기상청',
      warn='영국 자료는 지금 수집만 하고 화면에 올리지 않았습니다 — '
           'Met Office 무료 플랜의 재배포 조항 확인이 끝나기 전에는 쓰지 않습니다.'),

 dict(kind='table', n='2', title='네 나라에서 모으고 있습니다',
      head=['나라', '수집기', '무엇을 받는가'],
      rows=[['한국', '15개', '기상청 ASOS·AWS·동네예보·특보·낙뢰·해양·고층, 해양수산부 조위, 산림청 산불'],
            ['일본', '3개',  '일본 기상청 AMeDAS 실황, 기상경보, 지명'],
            ['영국', '1개',  '영국 기상청(Met Office) DataHub 지점예보 48시간'],
            ['미국', '7개',  'NOAA(구름합성·부이·우주기상·METAR·쓰나미), NASA(위성타일·화재), USGS 지진'],
            ['국제', '나머지','ECMWF, Copernicus CAMS, GDACS, GDELT, Celestrak, Open-Meteo']],
      kicker='한 나라만 보면 그 나라 모델이 옳은지 알 수 없습니다.'),

 dict(kind='table', n='3', title='"보여준다"가 무슨 뜻인가',
      lead='저희가 파는 것은 자료가 아니라 자료를 보이게 만드는 방법입니다.',
      head=['방식', '개수', '무엇을 그리나', '지금 쓰는 예'],
      rows=[['격자', '19', '값을 색으로 칠한 면', '기온·풍속·파고·대기질'],
            ['위성영상', '13', '영상을 지구 표면에 입힘', '천리안2A·히마와리·수오미'],
            ['지점', '12', '관측소를 점으로', 'ASOS·AMeDAS·부이·조위'],
            ['사건', '8', '일어난 일을 표시', '지진·산불·낙뢰·재난경보'],
            ['경로', '3', '움직인 자취', '태풍 진로·철새·바다거북'],
            ['바람', '1', '입자 애니메이션', '전지구 바람 흐름'],
            ['관측범위', '2', '어디까지 봤나', '위성 관측 영역']],
      kicker='위경도와 시각만 있으면 무엇이든 지구 위에 올릴 수 있습니다.'),

 dict(kind='duo', n='4', title='지구본 위 — 실제 화면',
      items=[('globe-gk2a.png',  '위성영상',
              '천리안2A 적외 11.2µm · 70~190°E\n밤에도 구름이 보입니다'),
             ('globe-wave.png',  '격자 — 파고',
              '태풍 주변에서 파고가 오르고\n너울이 퍼져 나가는 모습')]),

 dict(kind='duo', n='4', title='지구본 위 — 실제 화면',
      items=[('globe-hazard.png','사건',
              '산불(주황 점)·지진·쓰나미를 한 화면에\n오늘 시베리아 산불이 잡혀 있습니다'),
             ('globe-cyclone.png','경로 + 위성',
              '천리안 영상 위에 태풍과 예상경로\n기관별 경로를 겹쳐서 그대로')]),

 dict(kind='shot', n='5', title='차트로도 풉니다',
      img='page-verify-crop.png',
      caps=['/verify — GFS와 ECMWF 가 24시간 전에 내놓은 값을 기상청 ASOS 97지점 관측과 대조',
            '/research — 지점·기간·변수를 골라 시계열과 분포를 뽑는 화면',
            '/station — 관측소 한 곳을 깊게 보는 화면'],
      warn='화면의 집계는 현재 1일치입니다. 원본 아카이브는 7월 26일부터 쌓고 있습니다.'),

 dict(kind='stats', n='6', title='핵심 자산 — 예보 검증 아카이브',
      note='2026년 7월 26일부터 매시간 저장',
      stats=[('97', '곳', '기상청 관측지점'),
             ('2', '종', '모델\nGFS · ECMWF'),
             ('72', '시간', '예보 길이'),
             ('11', '일', '지금까지 쌓인 기간')],
      body='예보는 발표되고 사라집니다. 새 예보가 나오면 이전 예보는 덮어써집니다.\n'
           '모델 정확도를 지역별·조건별로 따지려면 그 순간에 박제해 두는 수밖에 없습니다.\n'
           '이 자산은 돈으로 살 수 없고 시간으로만 만들어집니다.',
      warn='솔직히 적습니다. 현재 11일치입니다. 얇습니다. '
           '다만 오늘 똑같이 시작하는 곳은 오늘부터 0일입니다.'),

 dict(kind='honest', n='7', title='없는 것 — 먼저 밝힙니다',
      lead='밋업에서 드러날 것을 제안서에서 먼저 말하는 편이 서로의 시간을 아낍니다.\n'
           'earthus 는 서비스 화면에서도 같은 원칙을 씁니다.',
      items=['학습된 AI 모델이 없습니다 — 코드 전수 확인 결과 TensorFlow·PyTorch·scikit-learn 사용 0건',
             '자체 위성·레이더가 없습니다 — 공개 자료를 받아 처리합니다',
             '전력·발전 도메인 지식이 없습니다',
             '예보를 직접 생산하지 않습니다 — 기관 예보를 모으고, 맞았는지 검증합니다'],
      kicker='저희가 내놓는 것은 모델이 아니라,\n모델을 만들고 채점하고 눈으로 확인하는 바탕입니다.'),

 dict(kind='offer', n='8', title='SK텔레콤', tag='AI 기반기술 연계 · 빅데이터',
      problem='AI 모델은 학습 데이터보다 평가와 설명에서 먼저 막힙니다.\n'
              '"이게 기존 수치예보보다 정말 나은가?" 에 답하려면\n'
              '같은 시각·같은 지점의 예보와 실측이 쌍으로 있어야 합니다.',
      head=['무엇', '지구본 위', '차트로'],
      rows=[['모델 채점판', '지점별 오차를 색으로 칠한 지도', '모델별·예보시간별 오차 곡선'],
            ['SKT 모델 vs 수치예보', '같은 시각·같은 지점을 나란히', '승패 누적, 조건별 우열'],
            ['오차의 지리적 패턴', '해안·산지·도심을 따라 몰리는 모습', '지형·시간대별 분해']],
      kicker='표에서 "RMSE 1.8" 은 숫자입니다.\n지구 위에서는 "동해안을 따라 붉게 몰려 있다" 가 됩니다.'),

 dict(kind='offer', n='9', title='LS ELECTRIC', tag='기후대응 / 신재생에너지 · VPP',
      problem='VPP 는 발전량을 예측해 전력시장에 입찰합니다.\n'
              '예측이 틀리면 정산에서 그대로 손실입니다.\n'
              '그런데 "내 예보가 어디서 얼마나 틀리는가" 를 아는 사업자는 드뭅니다.',
      head=['무엇', '지구본 위', '차트로'],
      rows=[['발전소 지도', '발전소 위치에 예보 신뢰도를 색으로', '사이트별 오차 이력'],
            ['오늘의 위험', '오차가 커질 조건이 다가오는 모습', '시간대별 신뢰구간'],
            ['모델 가중치', '지역별로 어느 모델을 믿어야 하는지', 'GFS vs ECMWF 조건별 승률'],
            ['일사·바람 흐름', '구름이 발전소 위를 지나가는 영상', '발전량 예측 대비 실적']],
      kicker='입찰 담당자는 지금 예보 숫자 하나를 봅니다.\n'
             '이 화면은 거기에 "이 지점 이 시간대에는 평소 이만큼 틀린다" 를 붙여 줍니다.'),

 dict(kind='offer', n='10', title='두산에너빌리티', tag='해상풍력 · AI/빅데이터',
      problem='해상풍력은 O&M 출항 판단이 곧 비용입니다.\n'
              '파고·풍속 예보가 틀리면 배를 띄웠다 돌아오거나,\n'
              '띄울 수 있었던 날을 놓칩니다.',
      head=['무엇', '지구본 위', '차트로'],
      rows=[['출항 창', '단지 주변 파고·풍속을 시간축으로 재생', '5일 중 나갈 수 있는 시간대'],
            ['태풍 판단', '기관 4곳 경로를 겹쳐서 그대로 표시', '기관별 강도·도달시각 차이'],
            ['예보 신뢰도', '해상 지점별 과거 오차폭', '파고 예보가 언제 빗나가는가'],
            ['실측 대조', '인근 부이 실측을 예보 위에 겹침', '예보-실측 편차 이력']],
      warn='태풍 경로는 하나로 합치지 않습니다. 기관이 갈릴 때가 정확히 판단이 어려운 때이고, '
           '그 갈림 자체가 정보이기 때문입니다. 이 화면은 이미 서비스에 있습니다.'),

 dict(kind='flow', n='11', title='지구를 쓰면 여기서 더 갈 수 있습니다',
      lead='표는 값을 보여주고, 지구는 관계를 보여줍니다.',
      pairs=[('엑셀에서 오차 1.8 은 그냥 1.8', '지구 위에서는 "산맥을 넘는 순간 커진다"'),
             ('발전소 200곳의 실적표는 200줄', '지구 위에서는 "이 구름대가 지나간 곳만 떨어졌다"'),
             ('태풍 경로 4개는 좌표 목록', '지구 위에서는 "어느 시점부터 갈라지는지"')],
      steps=['보이게 만든다', '사람이 빨리 판단한다', '규칙으로 옮긴다', '알고리즘이 된다'],
      foot='저희가 잘하는 것은 첫 칸입니다. 협력사가 잘하는 것은 뒤 칸입니다.\n'
           '첫 칸이 없으면 뒤 칸은 감으로 시작합니다.'),

 dict(kind='table', n='12', title='PoC 제안 — 3개월',
      head=['시기', '할 일'],
      rows=[['1개월', '대상 지점·자산 확정, 검증 지표 합의 (무엇을 "맞았다"고 볼 것인가)'],
            ['2개월', '해당 지점 아카이브 확장 + 지구본 시연 화면 1차'],
            ['3개월', '차트·리포트 완성 + 실제 의사결정 1건에 적용해 결과 보고']],
      warn='3개월 뒤 아카이브는 약 100일치입니다. 계절 하나를 겨우 넘깁니다. '
           '1년 단위 분석이 필요하면 그 시점에 착수하는 것이 맞고, 지금 된다고 말하지 않겠습니다.'),

 dict(kind='honest', n='13', title='저희가 일하는 방식',
      lead='협업 상대를 고르실 때 도움이 될 것 같아 적습니다.',
      items=['없는 것을 있다고 쓰지 않습니다 — 화면에도 "항공기·선박 … 지금은 아직 없습니다" 라고 적혀 있습니다',
             '출처·표본수·기준을 값과 함께 보여줍니다 — 숫자만 크게 띄우지 않습니다',
             '예보하지 않습니다 — 기관 예보를 옮기고, 맞았는지 확인해 드립니다',
             '라이선스가 확인되지 않은 자료는 쓰지 않습니다 — 영국 자료를 수집만 하는 이유입니다',
             '안전 정보에는 값을 매기지 않습니다 — 태풍·지진·쓰나미는 유료화 대상에서 제외했습니다']),

 dict(kind='closing', n='14', title='저희가 가려는 곳',
      body='earthus 는 기상정보만 보여주는 서비스로 남을 생각이 없습니다.\n\n'
           '지금 하는 일은 바탕을 쌓는 일입니다. 네 나라 기상청 자료를 매시간 모으고,\n'
           '사라지는 예보를 박제하고, 그것을 지구 위에 얹어 사람이 볼 수 있게 만드는 일입니다.\n\n'
           '그 위에 기상 AI 를 올리는 것을 고민하고 있습니다. 지금은 모델이 없습니다.\n'
           '다만 순서를 그렇게 잡았습니다 — 먼저 채점표를 만들고, 그 다음에 선수를 키웁니다.',
      cards=[('산업', '발전·해운·건설·농업처럼\n하늘과 바다에 매출이 걸린 곳에서,\n숫자가 아니라 화면으로 판단하게'),
             ('교육', '아이가 태풍이 왜 생기는지,\n왜 기관마다 다르게 말하는지를\n글이 아니라 돌려 보면서 알게')],
      kicker='각 산업과 교육이 함께 쓰는 한국형 3D 지구 서비스',
      img='globe-cyclone.png'),

 # ⚠️ 빈 칸은 PD 가 채운다. config.local.js 의 BUSINESS 가 비어 있어 넣을 값이 없다.
 #    앱 하단 사업자 고지에도 같은 값이 필요하니 한 번에 채우면 두 곳이 해결된다.
 dict(kind='table', n='15', title='회사 정보',
      head=['항목', '내용'],
      rows=[['상호', '(작성 필요)'],
            ['대표자', '(작성 필요)'],
            ['사업자등록번호', '(작성 필요)'],
            ['주소', '(작성 필요)'],
            ['서비스', 'https://earthus.net  ·  2026년 8월 4일 정식 오픈'],
            ['연락처', '(작성 필요)']],
      kicker='본 제안서의 모든 수치는 2026년 8월 6일 기준 실측값입니다.'),
]


# ══════════════════════════════════════════════════════ PPTX
def build():
    from pptx import Presentation
    from pptx.util import Inches as In, Pt, Emu
    from pptx.dml.color import RGBColor as C
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    def rgb(h): return C(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    prs = Presentation()
    prs.slide_width, prs.slide_height = In(W), In(H)
    blank = prs.slide_layouts[6]

    def new(dark=True):
        sl = prs.slides.add_slide(blank)
        f = sl.background.fill; f.solid(); f.fore_color.rgb = rgb(INK if dark else PAPER)
        return sl

    def tb(sl, x, y, w, h, mid=False):
        t = sl.shapes.add_textbox(In(x), In(y), In(w), In(h)).text_frame
        t.word_wrap = True
        t.margin_left = t.margin_right = t.margin_top = t.margin_bottom = 0
        if mid: t.vertical_anchor = MSO_ANCHOR.MIDDLE
        return t

    def line(tf, text, size, color, bold=False, first=False, sp_before=0, sp_after=0,
             lh=1.35, align=None, font='Calibri'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.line_spacing = lh
        p.space_before = Pt(sp_before); p.space_after = Pt(sp_after)
        if align: p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = rgb(color); r.font.name = font
        return p

    def rect(sl, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
        s = sl.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
        s.line.fill.background(); s.shadow.inherit = False
        return s

    def chapter(sl, n, title):
        """번호를 청록 원에 넣는다 — 지구를 닮은 이 덱의 되풀이 요소."""
        if n:
            c = sl.shapes.add_shape(MSO_SHAPE.OVAL, In(M), In(0.52), In(0.46), In(0.46))
            c.fill.background(); c.line.color.rgb = rgb(TEAL); c.line.width = Pt(1.2)
            c.shadow.inherit = False
            t = c.text_frame; t.margin_left = t.margin_right = 0
            t.vertical_anchor = MSO_ANCHOR.MIDDLE
            line(t, n, 13, TEAL, bold=True, first=True, align=PP_ALIGN.CENTER)
            tx = M + 0.68
        else:
            tx = M
        t = tb(sl, tx, 0.5, W - tx - M, 0.55)
        line(t, title, 26, PAPER, bold=True, first=True)

    def warnbox(sl, y, text, x=None, w=None):
        """⚠️ x·w 를 받아야 한다. 예전엔 항상 M 에 그린 뒤 도형을 옮기는 꼼수를 썼는데,
           바탕만 옮겨지고 글자는 제자리에 남아 **주황 빈 상자**가 나왔다 (실제로 겪음)."""
        x = M if x is None else x
        w = w or (W - M*2)
        per = int(w * 12.4)                    # 폭에 들어가는 글자 수 어림
        h = 0.34 + 0.215 * (len(text) // per + 1)
        rect(sl, x, y, w, h, '1E1710')
        rect(sl, x, y, 0.05, h, AMBER)
        t = tb(sl, x + 0.28, y + 0.15, w - 0.5, h - 0.28)
        line(t, '\u26a0  ' + text, 10.5, 'E8C79A', first=True, lh=1.5)
        return y + h

    # ── 슬라이드별 ────────────────────────────────────────────
    for s in SLIDES:
        k = s['kind']

        if k == 'cover':
            sl = new()
            p = os.path.join(IMG, s['img'])
            if os.path.exists(p):
                sl.shapes.add_picture(p, In(4.9), In(0), height=In(H))
            rect(sl, 0, 0, 6.4, H, INK)          # 왼쪽 글자 자리를 덮어 대비를 만든다
            t = tb(sl, M, 2.05, 5.2, 3.4)
            line(t, s['sub'], 12, TEAL, bold=True, first=True, sp_after=14)
            line(t, s['title'], 40, PAPER, bold=True, sp_after=16, lh=1.15)
            line(t, s['sub2'], 11.5, MUTE, lh=1.7)

        elif k == 'thesis':
            sl = new()
            t = tb(sl, M, 0, W - M*2, H, mid=True)
            line(t, s['title'], 36, PAPER, bold=True, first=True, sp_after=22, lh=1.25)
            line(t, s['body'], 18, TEAL, lh=1.6, sp_after=26)
            line(t, s['foot'], 12, MUTE)

        elif k == 'stats':
            sl = new(); chapter(sl, s['n'], s['title'])
            t = tb(sl, M, 1.14, W - M*2, 0.3)
            line(t, s['note'], 11, MUTE, first=True)
            n = len(s['stats']); gap = 0.3
            cw = (W - M*2 - gap*(n-1)) / n
            for i, (v, u, lab) in enumerate(s['stats']):
                x = M + i*(cw+gap)
                rect(sl, x, 1.62, cw, 2.05, INK2)
                t = tb(sl, x+0.32, 1.92, cw-0.6, 0.72)
                p = t.paragraphs[0]; p.line_spacing = 1.0
                r = p.add_run(); r.text = v
                r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = rgb(TEAL); r.font.name='Calibri'
                r2 = p.add_run(); r2.text = ' ' + u
                r2.font.size = Pt(15); r2.font.color.rgb = rgb(MUTE); r2.font.name='Calibri'
                t2 = tb(sl, x+0.32, 2.72, cw-0.6, 0.7)
                line(t2, lab, 10.5, PAPER, first=True, lh=1.45)
            y = 4.02
            if s.get('body'):
                nl = s['body'].count('\n') + 1
                t = tb(sl, M, y, W - M*2, 0.3*nl + 0.2)
                line(t, s['body'], 12, PAPER, first=True, lh=1.72)
                y += 0.3*nl + 0.34
            if s.get('warn'): warnbox(sl, y, s['warn'])

        elif k == 'table':
            sl = new(); chapter(sl, s['n'], s['title'])
            y = 1.15
            if s.get('lead'):
                t = tb(sl, M, y, W - M*2, 0.32)
                line(t, s['lead'], 12, MUTE, first=True); y += 0.5
            head, rows = s['head'], s['rows']
            nc = len(head)
            widths = ([1.5, 1.2, W-M*2-2.7] if nc == 3 and len(head[2]) > 6
                      else [1.4, 0.9, 3.5, W-M*2-5.8] if nc == 4
                      else [1.5, W-M*2-1.5] if nc == 2
                      else [(W-M*2)/nc]*nc)
            rect(sl, M, y, W-M*2, 0.5, '17233A')
            cx = M
            for i, hcell in enumerate(head):
                t = tb(sl, cx+0.22, y+0.14, widths[i]-0.3, 0.26)
                line(t, hcell, 10.5, TEAL, bold=True, first=True)
                cx += widths[i]
            y += 0.5
            for ri, row in enumerate(rows):
                rh = 0.62 if max(len(c) for c in row) < 56 else 0.74
                if ri % 2 == 0: rect(sl, M, y, W-M*2, rh, '0E1828')
                cx = M
                for i, cell in enumerate(row):
                    t = tb(sl, cx+0.22, y+0.19, widths[i]-0.3, rh-0.28)
                    line(t, cell, 10.5, PAPER if i else TEAL,
                         bold=(i == 0), first=True, lh=1.4)
                    cx += widths[i]
                y += rh
            y += 0.26
            if s.get('kicker'):
                t = tb(sl, M, y, W-M*2, 0.4)
                line(t, s['kicker'], 15, TEAL, bold=True, first=True); y += 0.55
            if s.get('warn'): warnbox(sl, y, s['warn'])

        elif k == 'duo':
            sl = new(); chapter(sl, s['n'], s['title'])
            gap = 0.45; cw = (W - M*2 - gap) / 2
            # ⚠️ 그림 높이를 비율로 직접 구한다. 안 그러면 캡션이 그림 위에 얹힌다 (실제로 겪음).
            from PIL import Image as PILImage
            ih = 0
            for fn, _, _ in s['items']:
                p = os.path.join(IMG, fn)
                if os.path.exists(p):
                    iw, ihpx = PILImage.open(p).size
                    ih = max(ih, cw * ihpx / iw)
            top = 1.32
            for i, (fn, cap, desc) in enumerate(s['items']):
                x = M + i*(cw+gap)
                p = os.path.join(IMG, fn)
                if os.path.exists(p):
                    sl.shapes.add_picture(p, In(x), In(top), width=In(cw))
                cy = top + ih + 0.24
                t = tb(sl, x, cy, cw, 0.32)
                line(t, cap, 14, TEAL, bold=True, first=True)
                t = tb(sl, x, cy + 0.38, cw, 0.86)
                line(t, desc, 11, PAPER, first=True, lh=1.55)

        elif k == 'shot':
            sl = new(); chapter(sl, s['n'], s['title'])
            from PIL import Image as PILImage
            p = os.path.join(IMG, s['img'])
            iw_in = 5.75
            if os.path.exists(p):
                px, py = PILImage.open(p).size
                sl.shapes.add_picture(p, In(M), In(1.25), width=In(iw_in))
            x = M + iw_in + 0.55
            cw = W - x - M
            t = tb(sl, x, 1.3, cw, 3.2)
            for i, c in enumerate(s['caps']):
                head, _, rest = c.partition(' — ')
                line(t, head, 13, TEAL, bold=True, first=(i == 0), sp_before=0 if i == 0 else 15)
                line(t, rest, 11, PAPER, lh=1.55)
            if s.get('warn'):
                warnbox(sl, 4.7, s['warn'], x=x, w=cw)

        elif k == 'honest':
            sl = new(); chapter(sl, s['n'], s['title'])
            t = tb(sl, M, 1.16, W - M*2, 0.72)
            line(t, s['lead'], 12, MUTE, first=True, lh=1.6)
            y = 2.05
            for it in s['items']:
                rect(sl, M, y, W-M*2, 0.62, INK2)
                d = sl.shapes.add_shape(MSO_SHAPE.OVAL, In(M+0.3), In(y+0.26), In(0.11), In(0.11))
                d.fill.solid(); d.fill.fore_color.rgb = rgb(AMBER)
                d.line.fill.background(); d.shadow.inherit = False
                head, _, rest = it.partition(' — ')
                t = tb(sl, M+0.62, y+0.16, W-M*2-0.9, 0.4)
                p = t.paragraphs[0]; p.line_spacing = 1.3
                r = p.add_run(); r.text = head
                r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = rgb(PAPER); r.font.name='Calibri'
                if rest:
                    r2 = p.add_run(); r2.text = '   ' + rest
                    r2.font.size = Pt(10.5); r2.font.color.rgb = rgb(MUTE); r2.font.name='Calibri'
                y += 0.74
            if s.get('kicker'):
                t = tb(sl, M, y+0.16, W-M*2, 0.8)
                line(t, s['kicker'], 15, TEAL, bold=True, first=True, lh=1.5)

        elif k == 'offer':
            sl = new(); chapter(sl, s['n'], s['title'] + '   ' + s['tag'])
            t = tb(sl, M, 1.16, W - M*2, 0.95)
            line(t, s['problem'], 12.5, PAPER, first=True, lh=1.62)
            y = 2.3
            head, rows = s['head'], s['rows']
            widths = [2.5, 4.4, W-M*2-6.9]
            rect(sl, M, y, W-M*2, 0.4, '17233A')
            cx = M
            for i, hc in enumerate(head):
                t = tb(sl, cx+0.22, y+0.09, widths[i]-0.3, 0.24)
                line(t, hc, 10.5, TEAL, bold=True, first=True)
                cx += widths[i]
            y += 0.4
            for ri, row in enumerate(rows):
                if ri % 2 == 0: rect(sl, M, y, W-M*2, 0.5, '0E1828')
                cx = M
                for i, cell in enumerate(row):
                    t = tb(sl, cx+0.22, y+0.14, widths[i]-0.3, 0.3)
                    line(t, cell, 10.5, PAPER if i else TEAL, bold=(i == 0), first=True)
                    cx += widths[i]
                y += 0.5
            y += 0.28
            if s.get('kicker'):
                t = tb(sl, M, y, W-M*2, 0.8)
                line(t, s['kicker'], 14, TEAL, bold=True, first=True, lh=1.55)
            if s.get('warn'): warnbox(sl, y, s['warn'])

        elif k == 'flow':
            sl = new(); chapter(sl, s['n'], s['title'])
            t = tb(sl, M, 1.16, W-M*2, 0.34)
            line(t, s['lead'], 15, TEAL, bold=True, first=True)
            y = 1.75
            for a, b in s['pairs']:
                rect(sl, M, y, 5.1, 0.62, '0E1828')
                t = tb(sl, M+0.28, y+0.19, 4.6, 0.3)
                line(t, a, 11, MUTE, first=True)
                t = tb(sl, M+5.35, y+0.19, 0.4, 0.3)
                line(t, '→', 13, TEAL, bold=True, first=True)
                rect(sl, M+5.85, y, W-M*2-5.85, 0.62, INK2)
                t = tb(sl, M+6.13, y+0.19, W-M*2-6.3, 0.3)
                line(t, b, 11, PAPER, bold=True, first=True)
                y += 0.76
            y += 0.3
            n = len(s['steps']); gap = 0.3
            cw = (W - M*2 - gap*(n-1)) / n
            for i, st in enumerate(s['steps']):
                x = M + i*(cw+gap)
                rect(sl, x, y, cw, 0.66, '17233A')
                t = tb(sl, x, y+0.21, cw, 0.3)
                line(t, st, 12, TEAL if i < n-1 else AMBER, bold=True, first=True,
                     align=PP_ALIGN.CENTER)
                if i < n-1:
                    t = tb(sl, x+cw+0.03, y+0.21, 0.24, 0.3)
                    line(t, '›', 15, MUTE, first=True, align=PP_ALIGN.CENTER)
            y += 0.95
            t = tb(sl, M, y, W-M*2, 0.7)
            line(t, s['foot'], 11.5, MUTE, first=True, lh=1.6)

        elif k == 'closing':
            sl = new()
            p = os.path.join(IMG, s['img'])
            if os.path.exists(p):
                pic = sl.shapes.add_picture(p, In(7.2), In(0), height=In(H))
            rect(sl, 0, 0, 7.6, H, INK)
            chapter(sl, s['n'], s['title'])
            t = tb(sl, M, 1.3, 6.3, 2.5)
            line(t, s['body'], 11.5, PAPER, first=True, lh=1.72)
            y = 4.1
            for i, (h_, d) in enumerate(s['cards']):
                x = M + i*3.25
                rect(sl, x, y, 3.0, 1.5, INK2)
                t = tb(sl, x+0.26, y+0.2, 2.5, 0.3)
                line(t, h_, 13, TEAL, bold=True, first=True)
                t = tb(sl, x+0.26, y+0.56, 2.5, 0.85)
                line(t, d, 10, PAPER, first=True, lh=1.5)
            t = tb(sl, M, 5.95, 6.3, 0.7)
            line(t, s['kicker'], 15, AMBER, bold=True, first=True, lh=1.45)

    prs.save(PPTX)
    return len(SLIDES)


if __name__ == '__main__':
    n = build()
    kb = os.path.getsize(PPTX) / 1024
    print(f'PPTX {kb:,.0f} KB · {n}장  →  {os.path.basename(PPTX)}')
