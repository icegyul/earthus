#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""만들어진 PPTX 를 **다시 읽어** 그대로 HTML 로 그린다 (눈 검수용).

⚠️ 왜 이렇게 하나
   이 기계에 LibreOffice·pdftoppm 이 없어 PPTX 를 그림으로 못 바꾼다.
   그렇다고 원본 자료(SLIDES)로 미리보기를 따로 그리면 **PPTX 와 어긋날 수 있다** —
   미리보기는 멀쩡한데 실제 파일이 깨져 있어도 모른다.
   그래서 python-pptx 로 **완성된 파일을 다시 읽어** 도형 위치·글자를 그대로 옮긴다.
   글자가 넘치거나 겹치면 여기서 보인다.

⚠️ 폰트는 파워포인트가 그리는 것이라 미리보기와 완전히 같지는 않다.
   위치·크기·겹침을 보는 용도다.

쓰는 법:  python3 docs/proposals/deck_preview.py
"""
import base64, io, os, subprocess

from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(HERE, 'earthus-협업제안서-2026.pptx')
OUT  = os.path.join(HERE, 'deck-preview.html')
PX   = 96.0 / 914400.0        # EMU → px (96dpi)


def esc(s):
    return (s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;'))


def color_of(fill, default=None):
    try:
        if fill.type is not None and fill.fore_color.type is not None:
            return '#%02X%02X%02X' % tuple(fill.fore_color.rgb)
    except Exception:
        pass
    return default


def main():
    prs = Presentation(PPTX)
    SW, SH = prs.slide_width * PX, prs.slide_height * PX
    out = [f'''<meta charset="utf-8"><style>
body{{margin:0;background:#31384a;font-family:"Apple SD Gothic Neo",sans-serif}}
.wrap{{padding:18px}}
.s{{position:relative;width:{SW}px;height:{SH}px;margin:0 auto 18px;overflow:hidden;
    box-shadow:0 6px 24px rgba(0,0,0,.5)}}
.n{{color:#aab;font:12px monospace;margin:0 auto 4px;width:{SW}px}}
.t{{position:absolute;white-space:pre-wrap;overflow-wrap:break-word}}
.t p{{margin:0}}
</style><div class="wrap">''']

    for i, sl in enumerate(prs.slides, 1):
        bg = color_of(sl.background.fill, '#101010')
        out.append(f'<div class="n">slide {i}</div><div class="s" style="background:{bg}">')
        for sh in sl.shapes:
            L, T = sh.left * PX, sh.top * PX
            Wd, Ht = sh.width * PX, sh.height * PX
            if sh.shape_type == 13 or sh.__class__.__name__ == 'Picture':   # PICTURE
                b64 = base64.b64encode(sh.image.blob).decode()
                out.append(f'<img src="data:{sh.image.content_type};base64,{b64}" '
                           f'style="position:absolute;left:{L}px;top:{T}px;'
                           f'width:{Wd}px;height:{Ht}px;object-fit:cover">')
                continue
            if sh.has_text_frame is False or not sh.text_frame.text.strip():
                c = color_of(sh.fill)
                ln = None
                try:
                    ln = '#%02X%02X%02X' % tuple(sh.line.color.rgb)
                except Exception:
                    pass
                if c or ln:
                    rad = '50%' if 'OVAL' in str(sh.shape_type) else '0'
                    st = f'background:{c};' if c else ''
                    st += f'border:1px solid {ln};' if ln else ''
                    out.append(f'<div style="position:absolute;left:{L}px;top:{T}px;'
                               f'width:{Wd}px;height:{Ht}px;border-radius:{rad};{st}"></div>')
                continue
            # 글자
            tf = sh.text_frame
            html = []
            for p in tf.paragraphs:
                # ⚠️ python-pptx PP_ALIGN: LEFT=1 CENTER=2 RIGHT=3.
                #    예전에 {1:'center'} 로 잘못 읽어 가운데 정렬이 오른쪽으로 보였다.
                al = {2: 'center', 3: 'right'}.get(
                    getattr(p.alignment, 'value', None), 'left')
                lh = p.line_spacing or 1.3
                runs = []
                for r in p.runs:
                    sz = (r.font.size.pt if r.font.size else 12)
                    col = '#FFFFFF'
                    try:
                        col = '#%02X%02X%02X' % tuple(r.font.color.rgb)
                    except Exception:
                        pass
                    bold = 'font-weight:700;' if r.font.bold else 'font-weight:400;'
                    runs.append(f'<span style="font-size:{sz*1.333:.1f}px;color:{col};{bold}">'
                                f'{esc(r.text)}</span>')
                sb = (p.space_before.pt if p.space_before else 0) * 1.333
                sa = (p.space_after.pt if p.space_after else 0) * 1.333
                html.append(f'<p style="text-align:{al};line-height:{lh};'
                            f'margin:{sb}px 0 {sa}px">{"".join(runs) or "&nbsp;"}</p>')
            mid = 'display:flex;flex-direction:column;justify-content:center;' \
                  if str(tf.vertical_anchor) .startswith('MIDDLE') else ''
            out.append(f'<div class="t" style="left:{L}px;top:{T}px;width:{Wd}px;'
                       f'height:{Ht}px;{mid}">{"".join(html)}</div>')
        out.append('</div>')

    out.append('</div>')
    io.open(OUT, 'w', encoding='utf-8').write(''.join(out))
    print(f'미리보기 {os.path.getsize(OUT)/1024:,.0f} KB · {len(prs.slides)}장 → {os.path.basename(OUT)}')


if __name__ == '__main__':
    main()
