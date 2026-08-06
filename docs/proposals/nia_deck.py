#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""'26년 공공데이터 활용기업 AI 연계 맞춤형 성장 지원사업 — 제출용 사업계획서 PPTX

원본 내용: nia_form.py 의 사업계획서 2-1~2-7 (이미 채워 넣은 것과 같은 사실·숫자)
쓰는 법:  python3 docs/proposals/nia_deck.py

⚠️ 이 파일이 하는 일 — deck.py(오픈이노베이션용)와 같은 방식이다.
   이 기계에 pptxgenjs·LibreOffice·pdftoppm 이 없어 python-pptx 로 만들고,
   **완성된 PPTX 를 다시 읽어** HTML 로 그려 눈으로 검수한다 (nia_deck_preview.py).

⚠️ 색은 제품이 아니라 이 심사에 맞춘 것이다. 심사위원은 공공기관 실무자다.
   오픈이노베이션 제안서의 우주 남색은 대기업 피칭용이었다.
   여기는 흰 바탕 + 다크네이비 텍스트로 관공서 문서에 가깝게 갔다 —
   이게 튀지 않아야 할 자리다.

⚠️ 화면캡처 — PD 가 직접 찍어 보낸 것 6장 + 기존 오픈이노베이션 제안서용 8장을 함께 쓴다.
   PD 가 보낸 것 중 '07-globe-overlap-bug' 는 뺐다. 태풍·부이·철새 선이 뒤엉켜
   레이어 겹침 버그처럼 보이는데, 실제 버그인지 확인이 안 된 상태로
   정부 제출자료에 넣을 수 없다. 검증 없이 화면을 골라 넣지 않는다.

⚠️ 지어내지 않는다 — nia_form.py 와 같은 규칙.
   재무·인력은 FILL 로 비워 둔다. 로드맵은 목표가 아니라 가정과 계산식.
"""
import os
from docx.shared import RGBColor  # 색상값만 재사용 (Pt/Cm 은 pptx 쪽에서 새로 씀)

HERE = os.path.dirname(os.path.abspath(__file__))
IMG  = os.path.join(HERE, 'img')
NIA  = os.path.join(IMG, 'nia')
PPTX = os.path.join(HERE, 'NIA-사업계획서-earthus.pptx')

# ── 색 — 관공서 문서 톤 ──────────────────────────────────────────
INK   = '15202B'   # 거의 검정 (본문)
NAVY  = '0F2647'   # 짙은 남색 — 표지·제목 (지배색)
BLUE  = '1F5C8B'   # 강조 (링크·소제목)
AMBER = 'B5620A'   # 경고
PAPER = 'FFFFFF'
CARD  = 'F4F6F9'
LINE  = 'D9E0E8'
MUTE  = '5C6B7A'

W, H = 13.333, 7.5
M = 0.85

FILL = '［작성 필요］'

SLIDES = [
 dict(kind='cover', title="'26년 공공데이터 활용기업 AI 연계 맞춤형 성장 지원사업",
      sub='사업계획서', sub2='달루어(earthus) · 대표 신효정 · 2026년 8월',
      img='nia/01-typhoon-globe.png'),

 dict(kind='overview', n='1', title='한눈에 보기',
      lead='예보는 발표되고 사라집니다. earthus 는 한·일·영·미 기상청 공공데이터를 '
           '매시간 자동 수집해 3D 지구본 위에 보여주고, 그 예보가 실제로 맞았는지를 '
           '시간 단위로 보존해 공개합니다.',
      stats=[('2020', '설립', '개인사업자 · 정보통신업'),
             ('2026.08.04', '정식 오픈', 'earthus.net'),
             ('18', '개 기관', '4개국 공공데이터 연계'),
             ('54', '개', 'AWS Lambda 자동 수집 파이프라인')]),

 dict(kind='table', n='1-1', title='기업 개요',
      head=['항목', '내용'],
      rows=[['회사명', '달루어'],
            ['대표자', '신효정'],
            ['사업자등록번호', '419-03-01562'],
            ['설립일자', '2020. 01. 29.  (업력 6년 6개월 · 도약창업기업 구간)'],
            ['주소', '서울특별시 동작구 사당로16가길 102, 지층 좌측(사당동)'],
            ['기존 사업', '영상·미디어 콘텐츠 제작 (광고·영화·촬영장비 임대)'],
            ['신규 사업', 'earthus — 공공데이터 기반 3D 지구 기상·재난 시각화 서비스 (2026년 착수)']],
      widths=[3.2, 10.1], headrow=True,
      note='⚠️ 업력은 도약창업기업(4~7년) 구간이나, 공공데이터·AI 사업은 2026년에 '
           '착수한 신규 영역입니다. 기존 사업과 별개로 AI 도입 단계는 초기에 '
           '가깝다고 판단해, 도약 구간 추천 목록에 없는 「AI-Ready 구축」을 함께 신청합니다.'),

 dict(kind='shot2', n='2-1', title='사업 배경 — 예보는 발표되고 사라집니다',
      body='새 예보가 나오면 이전 예보는 덮어써지고, 지나고 나서 "그때 무엇이라고 '
           '했는지"를 되찾을 방법이 없습니다. 공공데이터포털·기상청 API허브는 '
           '방대한 자료를 공개하지만 대부분 "지금 값"만 제공하고 과거 발표분은 '
           '보관하지 않습니다. earthus 는 이 공백에서 출발했습니다 — 공공데이터를 '
           '받아 보여주는 데서 그치지 않고, 그 데이터가 실제로 맞았는지 함께 기록해 '
           '공개합니다.',
      imgs=['nia/01-typhoon-globe.png', 'page-verify-crop.png'],
      caps=['태풍 DOLPHIN — 천리안2A 위성영상 위에 기관별 예상경로를 겹쳐 표시',
            '예보 검증 화면 — GFS·ECMWF 예보를 기상청 실측과 대조']),

 dict(kind='shot2', n='2-2', title='서비스 소개 — 무엇을 보여주는가',
      body='웹 브라우저에서 3D 지구본을 돌려 보며 기상·해양·재난·천문 정보를 '
           '확인하는 서비스입니다. 설치가 필요 없고 한국어·영어를 지원하며, '
           '2026년 8월 4일 정식 오픈했습니다. 지진·산불·위성·철새 이동까지 '
           '공공데이터를 실시간으로 지구 위에 그립니다.',
      imgs=['nia/05-quake-detail.png', 'nia/04-wildfire.png'],
      caps=['지진 상세 — EMSC 관측을 규모·깊이·위치·시각과 함께 표시',
            '산불 — NASA VIIRS 위성영상으로 열입지 실시간 표시']),

 dict(kind='shot2', n='2-2b', title='서비스 소개 — 위성·생태 데이터까지',
      body='기상·재난에 그치지 않고 위성 궤도 카탈로그, 철새 이동 관측 같은 '
           '공공데이터도 함께 다룹니다. 다만 자료의 한계도 화면에 그대로 밝힙니다 — '
           '철새 이동 경로는 실제 비행 궤적이 아니라 출발지·도착지 두 지점을 이은 '
           '것이며, 화면에 그렇게 적어 둡니다. 없는 정밀도를 있는 것처럼 보이게 '
           '하지 않습니다.',
      imgs=['nia/03-satellite-detail.png', 'nia/06-migratory-bird.png'],
      caps=['위성 궤도 카탈로그 — Iridium 170 상세 정보',
            '철새 이동 관측 — 179건 · 14종 · 한계를 화면에 명시']),

 dict(kind='table', n='2-3', title='활용 중인 공공데이터',
      head=['제공기관', '데이터', '방식·주기'],
      rows=[['기상청 API허브', '지진통보·전지구관측(GTS)·낙뢰관측·평년값·태풍정보·산악기상',
             '매시간~'],
            ['공공데이터포털(환경공단)', '대기오염 실시간 측정정보, 측정소 정보', '매시간'],
            ['공공데이터포털(산림청)', '산불 발생 정보', '3시간'],
            ['공공데이터포털(해수부)', '조위관측 실시간, 이안류 지수', '매시간'],
            ['공공데이터포털(해양생물자원관)', '바다거북 이동경로, 바닷새 조사정보', '일 1회'],
            ['국립생태원 에코뱅크', '철새 이동 관측', '일 1회'],
            ['농림축산식품부', '철새도래지 정보', '일 1회']],
      widths=[3.5, 6.9, 1.9], headrow=True,
      note='기상청 API허브는 482개 API 활용 신청이 승인되어 있으며, 현재 실제 '
           '호출은 6종입니다. 승인된 나머지를 어떤 순서로 도입할지 컨설팅을 '
           '희망합니다.'),

 dict(kind='table', n='2-3b', title='시스템 구성 — 어떻게 처리하는가',
      head=['단계', '구성', '내용'],
      rows=[['① 수집', 'AWS Lambda 54개', '18개 기관을 매시간~6시간 주기로 자동 호출'],
            ['② 가공', 'Lambda + Python', '위성 원본→등경위도 변환, 관측값 격자 보간, 예보 시각별 분리'],
            ['③ 보관', 'Amazon S3', '2026-08-06 기준 누적 184MB · 일 10.2MB 증가'],
            ['④ 배포', 'CloudFront CDN', '전 세계 캐시 배포'],
            ['⑤ 표현', '브라우저 (Cesium WebGL)', '58개 레이어 · 격자·위성영상·지점·사건·경로·바람 7방식']],
      widths=[2.0, 3.6, 7.7], headrow=True,
      note='서버리스 구조로 고정비가 낮습니다. 2026-08-06 기준 서버 비용은 월 1만원 미만으로 추정됩니다.'),

 dict(kind='stats', n='3', title='핵심 자산 — 예보 검증 아카이브',
      note='2026년 7월 26일부터 매시간 저장',
      stats=[('97', '곳', '기상청 관측지점'),
             ('2', '종', '모델 GFS·ECMWF'),
             ('72', '시간', '예보 길이'),
             ('11', '일', '지금까지 쌓인 기간')],
      body='예보-실측 쌍은 시간이 지나야만 만들어지는 자산입니다. 지역별·조건별로 '
           '모델 정확도를 따지려면 발표되는 순간에 보존해 두는 수밖에 없습니다. '
           '이 기록은 향후 기상 AI 의 학습·평가에 쓰이는 정답지(ground truth)가 됩니다.',
      warn='솔직히 적습니다. 현재 11일치입니다. 얇습니다. 다만 오늘 똑같이 '
           '시작하는 곳은 오늘부터 0일이며, 저희는 매일 쌓고 있습니다.'),

 dict(kind='honest', n='4', title='없는 것 — 먼저 밝힙니다',
      lead='심사에서 확인될 것을 계획서에서 먼저 말하는 편이 정확한 평가에 도움이 '
           '된다고 생각합니다. earthus 는 서비스 화면에서도 같은 원칙을 씁니다.',
      items=['학습된 AI 모델이 없습니다 — 코드 전수 확인 결과 TensorFlow·PyTorch·scikit-learn 사용 0건',
             '전용 서버·자체 위성이 없습니다 — 서버리스 구조로 공개 자료를 받아 처리합니다',
             '재무·인력 실적은 신규 사업(2026년 착수) 특성상 아직 없습니다',
             '기상청 API허브 482개 중 실제 활용은 6종입니다 — 승인은 받았으나 아직 다 쓰지 못했습니다'],
      kicker='본 사업을 통해 채우고자 하는 것이 바로 이 빈칸입니다 — '
             '자료를 AI 학습 형태로 만드는 방법과, 다음에 무엇을 붙일지에 대한 방향입니다.'),

 dict(kind='program', n='5', title='희망 프로그램과 그 이유',
      rows=[('AI-Ready 구축', '공공데이터 수집·정제 / 구조화·라벨링 / 학습데이터 준비',
             '예보-실측 쌍을 매시간 쌓고 있으나, AI 학습에 쓸 수 있는 형태로 구조화·라벨링하는 '
             '방법을 정하지 못했습니다. 어떤 단위로 묶고 어떤 품질 기준으로 거를지 전문가 검토가 필요합니다.'),
            ('AI 서비스 고도화', '특허 출원 지원 / AI 모델 고도화',
             '예보를 시각 단위로 보존해 실측과 대조하는 방식이 핵심 자산인데 현재 아무런 권리 '
             '보호가 없습니다. 축적한 데이터 위에 AI 를 올리는 방향 설정도 함께 필요합니다.'),
            ('투자·해외 진출', '사업 체계화 및 투자 유치 / IR·네트워킹',
             '개인 구독만으로는 필요 데이터(항공기·선박 실시간 등) 구매 비용을 감당하기 어렵습니다.'),
            ('창업지원', '필요 공공데이터 탐색 및 매칭',
             '기상청 API허브 482개 중 실제 활용은 6종입니다. 레이더 합성 등 우선 도입할 데이터를 '
             '함께 골라 주시면 좋겠습니다.')]),

 dict(kind='table', n='6', title='비즈니스 모델',
      head=['구분', '내용', '상태'],
      rows=[['계속 무료', '태풍·지진·쓰나미 등 안전 정보, 구름·날씨, 예보 검증, 3D 학습 콘텐츠', '제공 중'],
            ['유료 (예정)', '위성 전체 목록·궤도 추적선, 관심 지점 알림 20곳', '개발 완료 · 판매 미개시'],
            ['유료 (준비 중)', '되감기·이력 조회, 항공기·선박 실시간 위치', '자료 축적·데이터 구매 필요']],
      widths=[2.4, 8.0, 3.1], headrow=True,
      note='⚠️ 안전 정보에는 값을 매기지 않는다는 원칙을 두고 있습니다. '
           '개인 구독 월 5,900원 / 연 49,000원 기준이나 최종 확정 전이며, '
           '사전등록 선착순 500명에게는 정가의 50%를 적용합니다.'),

 dict(kind='roadmap', n='7', title='로드맵',
      note='아래는 목표치가 아니라 가정과 계산식입니다.',
      rows=[('2026년 하반기', '검증 아카이브 100일치 축적, 변수 확대(강수·습도·일사)',
             '유료 기능 개시, 창립 멤버 500명 모집',
             '500명 × 연 24,500원 = 약 1,225만원 (전원 결제 가정 시)'),
            ('2027년', '검증 아카이브 1년치(계절 4개 확보), 모델 확대',
             '되감기·이력 개시, 산업용 API 제공 검토',
             '개인 구독 + 기업 데이터 제공'),
            ('2028년', '2년치 축적 — 전년 동기 비교 가능',
             '기상 AI 모델 학습·평가 착수, 해외 지역 확대',
             FILL)],
      warn='검증 아카이브는 2026-07-26 시작으로 현재 11일치입니다. 계절 단위 분석은 '
           '자료가 쌓이는 시간을 기다려야 하며, 앞당길 방법이 없습니다.'),

 dict(kind='honest', n='8', title='기대효과',
      lead='',
      items=['공공데이터 신뢰도의 가시화 — 값과 함께 그 값의 오차 이력을 공개적으로 보여줍니다',
             '기상 AI 학습·평가 기반 조성 — 예보-실측 쌍은 기상 AI 학습에 필요한 정답지입니다',
             '산업 활용 — 발전량 예측, 해상 작업 판단 등 예보 오차가 곧 비용인 분야에 활용 가능합니다',
             '교육 활용 — 태풍이 왜 생기는지, 왜 기관마다 다르게 말하는지 지구본으로 이해하게 합니다',
             '공공데이터 활용 저변 확대 — 개인 개발자도 서버리스 구조로 다기관 데이터를 통합할 수 있음을 보이는 사례입니다']),

 dict(kind='closing', title='공공데이터로 만드는 한국형 3D 지구 서비스',
      body='earthus 는 기상정보만 보여주는 서비스로 남을 생각이 없습니다.\n\n'
           '지금 하는 일은 바탕을 쌓는 일입니다. 18개 기관의 공공데이터를 매시간 '
           '모으고, 사라지는 예보를 보존하고, 그것을 지구 위에 얹어 사람이 볼 수 '
           '있게 만드는 일입니다.\n\n'
           '그 위에 AI를 올리는 것을 고민하고 있습니다. 지금은 모델이 없습니다. '
           '다만 순서를 그렇게 잡았습니다 — 먼저 채점표를 만들고, 그 다음에 '
           '선수를 키웁니다. 이번 지원사업이 그 다음 단계로 가는 발판이 되기를 '
           '바랍니다.',
      img='nia/01-typhoon-globe.png'),
]


def build():
    from pptx import Presentation
    from pptx.util import Inches as In, Pt
    from pptx.dml.color import RGBColor as C
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    def rgb(h): return C(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    prs = Presentation()
    prs.slide_width, prs.slide_height = In(W), In(H)
    blank = prs.slide_layouts[6]

    def new(bg=PAPER):
        sl = prs.slides.add_slide(blank)
        f = sl.background.fill; f.solid(); f.fore_color.rgb = rgb(bg)
        return sl

    def tb(sl, x, y, w, h, mid=False):
        t = sl.shapes.add_textbox(In(x), In(y), In(w), In(h)).text_frame
        t.word_wrap = True
        t.margin_left = t.margin_right = t.margin_top = t.margin_bottom = 0
        if mid: t.vertical_anchor = MSO_ANCHOR.MIDDLE
        return t

    def line(tf, text, size, color, bold=False, first=False, sp_before=0, sp_after=0,
             lh=1.35, align=None, font='Calibri'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.line_spacing = lh
        p.space_before = Pt(sp_before); p.space_after = Pt(sp_after)
        if align: p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = rgb(color); r.font.name = font
        return p

    def rect(sl, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, outline=None):
        s = sl.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
        if outline:
            s.line.color.rgb = rgb(outline); s.line.width = Pt(0.75)
        else:
            s.line.fill.background()
        s.shadow.inherit = False
        return s

    def pic_path(name):
        p = os.path.join(IMG, name)
        return p if os.path.exists(p) else None

    LOGO_B = pic_path('logo/earthus-lockup-black.png')
    LOGO_AR = None
    if LOGO_B:
        _w, _h = PILImage.open(LOGO_B).size
        LOGO_AR = _w / _h

    def chapter(sl, n, title, dark=False):
        col = PAPER if dark else NAVY
        if LOGO_B:
            lw = 1.15; lh_ = lw / LOGO_AR
            sl.shapes.add_picture(LOGO_B, In(W - M - lw), In(0.5 + (0.34-lh_)/2 if lh_<0.34 else 0.5),
                                   width=In(lw))
        if n:
            c = sl.shapes.add_shape(MSO_SHAPE.OVAL, In(M), In(0.5), In(0.5), In(0.5))
            c.fill.solid(); c.fill.fore_color.rgb = rgb(NAVY if not dark else PAPER)
            c.line.fill.background(); c.shadow.inherit = False
            t = c.text_frame; t.margin_left = t.margin_right = 0
            t.vertical_anchor = MSO_ANCHOR.MIDDLE
            line(t, n, 12, PAPER if not dark else NAVY, bold=True, first=True, align=PP_ALIGN.CENTER)
            tx = M + 0.68
        else:
            tx = M
        t = tb(sl, tx, 0.5, W - tx - M, 0.6)
        line(t, title, 22, col, bold=True, first=True)
        rect(sl, M, 1.14, W - M*2, 0.018, LINE if not dark else '3A4A5E')

    def warnbox(sl, y, text, x=None, w=None):
        x = M if x is None else x
        w = w or (W - M*2)
        per = int(w * 13)
        h = 0.3 + 0.2 * (len(text) // per + 1)
        rect(sl, x, y, w, h, '2A2015' if False else 'FDF3E8')
        rect(sl, x, y, 0.045, h, AMBER)
        t = tb(sl, x + 0.24, y + 0.13, w - 0.44, h - 0.24)
        line(t, '⚠  ' + text, 10, AMBER, first=True, lh=1.45)
        return y + h + 0.1

    def notebox(sl, y, text, x=None, w=None):
        x = M if x is None else x
        w = w or (W - M*2)
        per = int(w * 13)
        h = 0.3 + 0.2 * (len(text) // per + 1)
        t = tb(sl, x, y, w, h)
        line(t, '※ ' + text, 9.5, MUTE, first=True, lh=1.45)
        return y + h + 0.06

    def fit_img(sl, path, x, y, maxw, maxh):
        """비율 유지 · 상한 안에 맞춘다."""
        if not path: return maxh
        iw, ih = PILImage.open(path).size
        ar = iw / ih
        w, h = maxw, maxw / ar
        if h > maxh:
            h = maxh; w = maxh * ar
        sl.shapes.add_picture(path, In(x + (maxw-w)/2), In(y), width=In(w), height=In(h))
        return h

    for s in SLIDES:
        k = s['kind']

        if k == 'cover':
            sl = new(NAVY)
            p = pic_path(s['img'])
            if p:
                iw, ih = PILImage.open(p).size
                ar = iw/ih
                ph = H; pw = ph * ar
                sl.shapes.add_picture(p, In(W-pw), In(0), height=In(ph))
            rect(sl, 0, 0, 7.3, H, NAVY)
            lg = pic_path('logo/earthus-lockup-white.png')
            if lg:
                lw, lh_ = PILImage.open(lg).size
                pic_w = 2.9
                sl.shapes.add_picture(lg, In(M), In(1.7), width=In(pic_w))
                logo_bottom = 1.7 + pic_w * lh_ / lw
            else:
                logo_bottom = 1.7
            t = tb(sl, M, logo_bottom + 0.45, 6.0, 3.0)
            line(t, "공공데이터 활용기업 AI 연계 맞춤형 성장 지원사업", 11.5, '9FC1E0',
                 bold=True, first=True, sp_after=16, lh=1.4)
            line(t, s['sub'], 22, PAPER, bold=True, sp_after=18)
            line(t, s['sub2'], 11, 'C5D6E8', lh=1.7)

        elif k == 'overview':
            sl = new(); chapter(sl, s['n'], s['title'])
            t = tb(sl, M, 1.34, W-M*2, 0.9)
            line(t, s['lead'], 13, INK, first=True, lh=1.65)
            n = len(s['stats']); gap = 0.3
            cw = (W - M*2 - gap*(n-1)) / n
            y0 = 2.5
            for i, (v, u, lab) in enumerate(s['stats']):
                x = M + i*(cw+gap)
                rect(sl, x, y0, cw, 1.9, CARD, outline=LINE)
                t = tb(sl, x+0.24, y0+0.24, cw-0.48, 0.6)
                p = t.paragraphs[0]; p.line_spacing = 1.0
                r = p.add_run(); r.text = v
                r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = rgb(NAVY); r.font.name='Calibri'
                t2 = tb(sl, x+0.24, y0+0.78, cw-0.48, 0.34)
                line(t2, u, 12, BLUE, bold=True, first=True)
                t3 = tb(sl, x+0.24, y0+1.18, cw-0.48, 0.6)
                line(t3, lab, 9.5, MUTE, first=True, lh=1.4)

        elif k == 'table':
            sl = new(); chapter(sl, s['n'], s['title'])
            y = 1.32
            head, rows = s['head'], s['rows']
            widths = s.get('widths') or [(W-M*2)/len(head)]*len(head)
            rect(sl, M, y, W-M*2, 0.4, NAVY)
            cx = M
            for i, hc in enumerate(head):
                t = tb(sl, cx+0.2, y+0.08, widths[i]-0.3, 0.26)
                line(t, hc, 10.5, PAPER, bold=True, first=True)
                cx += widths[i]
            y += 0.4
            for ri, row in enumerate(rows):
                nlines = max((c.count('\n')+1) for c in row)
                rh = 0.4*nlines + 0.18
                if ri % 2 == 0: rect(sl, M, y, W-M*2, rh, CARD)
                cx = M
                for i, cell in enumerate(row):
                    t = tb(sl, cx+0.2, y+0.1, widths[i]-0.3, rh-0.18)
                    for j, ln_ in enumerate(cell.split('\n')):
                        line(t, ln_, 10, INK if i else NAVY, bold=(i==0), first=(j==0), lh=1.3,
                             sp_before=0 if j==0 else 1)
                    cx += widths[i]
                y += rh
            y += 0.18
            if s.get('note'): notebox(sl, y, s['note'])

        elif k == 'shot2':
            sl = new(); chapter(sl, s['n'], s['title'])
            t = tb(sl, M, 1.3, W-M*2, 0.95)
            line(t, s['body'], 11, INK, first=True, lh=1.6)
            gap = 0.4; cw = (W - M*2 - gap) / 2
            top = 2.5; maxh = 4.15
            for i, name in enumerate(s['imgs']):
                x = M + i*(cw+gap)
                p = pic_path(name)
                # ⚠️ 여기 rect(...,'FFFFFF00') 가 있었다. rgb() 는 앞 6자리만 읽어
                #    'FF','FF','FF' = 불투명 흰색이 되고, 방금 그린 그림을 그대로 덮었다.
                #    미리보기에서 그림이 통째로 안 보이는 원인이었다. 도형 자체를 지운다.
                h = fit_img(sl, p, x, top, cw, maxh)
                cy = top + h + 0.14
                t2 = tb(sl, x, cy, cw, 0.5)
                line(t2, s['caps'][i], 9.5, MUTE, first=True, lh=1.4)

        elif k == 'stats':
            sl = new(); chapter(sl, s['n'], s['title'])
            t = tb(sl, M, 1.24, W-M*2, 0.3)
            line(t, s['note'], 10.5, BLUE, bold=True, first=True)
            n = len(s['stats']); gap = 0.3
            cw = (W - M*2 - gap*(n-1)) / n
            y0 = 1.68
            for i, (v, u, lab) in enumerate(s['stats']):
                x = M + i*(cw+gap)
                rect(sl, x, y0, cw, 1.5, CARD, outline=LINE)
                t = tb(sl, x+0.24, y0+0.2, cw-0.48, 0.6)
                p = t.paragraphs[0]; p.line_spacing = 1.0
                r = p.add_run(); r.text = v
                r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = rgb(NAVY); r.font.name='Calibri'
                r2 = p.add_run(); r2.text = ' ' + u
                r2.font.size = Pt(12); r2.font.color.rgb = rgb(MUTE); r2.font.name='Calibri'
                t2 = tb(sl, x+0.24, y0+0.86, cw-0.48, 0.5)
                line(t2, lab, 9.5, INK, first=True, lh=1.35)
            y = y0 + 1.72
            if s.get('body'):
                nl = s['body'].count('\n') + 1
                t = tb(sl, M, y, W - M*2, 0.3*nl + 0.16)
                line(t, s['body'], 11.5, INK, first=True, lh=1.65)
                y += 0.3*nl + 0.3
            if s.get('warn'): warnbox(sl, y, s['warn'])

        elif k == 'honest':
            sl = new(); chapter(sl, s['n'], s['title'])
            y = 1.3
            if s.get('lead'):
                t = tb(sl, M, y, W-M*2, 0.6)
                line(t, s['lead'], 11, MUTE, first=True, lh=1.55)
                y += 0.66
            for it in s['items']:
                rect(sl, M, y, W-M*2, 0.56, CARD)
                d = sl.shapes.add_shape(MSO_SHAPE.OVAL, In(M+0.24), In(y+0.24), In(0.1), In(0.1))
                d.fill.solid(); d.fill.fore_color.rgb = rgb(BLUE)
                d.line.fill.background(); d.shadow.inherit = False
                head_, _, rest = it.partition(' — ')
                t = tb(sl, M+0.52, y+0.13, W-M*2-0.8, 0.34)
                p = t.paragraphs[0]; p.line_spacing = 1.25
                r = p.add_run(); r.text = head_
                r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = rgb(INK); r.font.name='Calibri'
                if rest:
                    r2 = p.add_run(); r2.text = '  —  ' + rest
                    r2.font.size = Pt(9.5); r2.font.color.rgb = rgb(MUTE); r2.font.name='Calibri'
                y += 0.64
            if s.get('kicker'):
                y += 0.14
                rect(sl, M, y, W-M*2, 0.78, NAVY)
                t = tb(sl, M+0.28, y+0.14, W-M*2-0.56, 0.5, mid=True)
                line(t, s['kicker'], 12.5, PAPER, bold=True, first=True, lh=1.45)

        elif k == 'program':
            sl = new(); chapter(sl, s['n'], s['title'])
            y = 1.32
            for stage, prog, why in s['rows']:
                rh = 1.16
                rect(sl, M, y, 2.6, rh, NAVY)
                t = tb(sl, M+0.18, y+0.16, 2.24, rh-0.3, mid=True)
                line(t, stage, 11.5, PAPER, bold=True, first=True, lh=1.3)
                rect(sl, M+2.6, y, W-M*2-2.6, rh, CARD)
                t = tb(sl, M+2.84, y+0.1, W-M*2-2.9, 0.32)
                line(t, prog, 10.5, BLUE, bold=True, first=True)
                t2 = tb(sl, M+2.84, y+0.42, W-M*2-2.9, rh-0.5)
                line(t2, why, 9.5, INK, first=True, lh=1.4)
                y += rh + 0.12

        elif k == 'roadmap':
            sl = new(); chapter(sl, s['n'], s['title'])
            t = tb(sl, M, 1.3, W-M*2, 0.3)
            line(t, s['note'], 10.5, BLUE, bold=True, first=True)
            y = 1.72
            head = ['시기', '데이터', '서비스', '수익 (가정)']
            widths = [1.9, 3.6, 3.5, 4.4]
            rect(sl, M, y, W-M*2, 0.4, NAVY)
            cx = M
            for i, hc in enumerate(head):
                t = tb(sl, cx+0.16, y+0.08, widths[i]-0.24, 0.26)
                line(t, hc, 10, PAPER, bold=True, first=True)
                cx += widths[i]
            y += 0.4
            for ri, row in enumerate(s['rows']):
                rh = 0.9
                if ri % 2 == 0: rect(sl, M, y, W-M*2, rh, CARD)
                cx = M
                for i, cell in enumerate(row):
                    t = tb(sl, cx+0.16, y+0.1, widths[i]-0.24, rh-0.18)
                    line(t, cell, 9.5, INK if i else NAVY, bold=(i==0), first=True, lh=1.35)
                    cx += widths[i]
                y += rh
            y += 0.16
            if s.get('warn'): warnbox(sl, y, s['warn'])

        elif k == 'closing':
            sl = new(NAVY)
            p = pic_path(s['img'])
            if p:
                iw, ih = PILImage.open(p).size
                ar = iw / ih
                # ⚠️ 오른쪽에 이미지를 채우고, 왼쪽은 불투명 남색 패널로 글자 대비를 확보한다.
                #    (이전 버전은 이미지를 그린 뒤 같은 자리를 다시 덮어 안 보이는 버그가 있었다)
                ph = H; pw = ph * ar
                sl.shapes.add_picture(p, In(W - pw), In(0), height=In(ph))
            rect(sl, 0, 0, 6.4, H, NAVY)
            t = tb(sl, M, 0.7, 5.5, H-1.4)
            line(t, s['title'], 22, PAPER, bold=True, first=True, sp_after=18, lh=1.3)
            line(t, s['body'], 11, 'D6E2EC', lh=1.75)

        else:
            raise ValueError(k)

    prs.save(PPTX)
    return len(SLIDES)


if __name__ == '__main__':
    n = build()
    print(f'PPTX {os.path.getsize(PPTX)/1024:,.0f} KB · {n}장 → {os.path.basename(PPTX)}')
