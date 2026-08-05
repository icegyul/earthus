# -*- coding: utf-8 -*-
"""
slides.py 하나에서 발표 자료 두 가지를 뽑는다.

  python3 build.py
    → index.html          (발표용 · 브라우저에서 ← → 로 넘김)
    → earthus-talk.pptx   (제출·인쇄용)

⚠️ 인쇄하면 어두운 배경에서 글이 안 보인다. 그래서 바탕은 밝게 간다.
   화면 캡처만 어두운 카드 위에 올려 대비를 준다.
"""

import io, os, sys, html, re

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
ROOT = os.path.abspath(os.path.join(HERE, "..", ".."))
LOGO = os.path.join(ROOT, "prototype", "logo-lockup.png")

sys.path.insert(0, HERE)
from slides import SLIDES  # noqa: E402


# ── 장 번호 자동 참조 ────────────────────────────────────────────────
# slides.py 의 «3부» → "3부(17–30장)", «장:제목조각» → "33장" 으로 바꾼다.
# ⚠️ 번호를 손으로 적으면 장 하나만 끼워 넣어도 전부 틀어진다. 그래서 여기서 센다.
def _resolve_refs(slides):
    titles = []            # (제목, 장번호) — 나온 순서대로
    sections = []          # (부 이름, 시작 장)
    for i, s in enumerate(slides, 1):
        titles.append((s.get("title", ""), i))
        if s["kind"] == "section":
            sections.append((s["part"], i))
    part_range = {}
    for idx, (part, start) in enumerate(sections):
        end = sections[idx + 1][1] - 1 if idx + 1 < len(sections) else len(slides)
        part_range[part] = "%s(%d–%d장)" % (part, start, end)

    def one(text):
        def sub_part(m):
            p = m.group(1)
            if p not in part_range:
                raise KeyError("모르는 부: " + p)
            return part_range[p]

        def sub_slide(m):
            frag = m.group(1)
            for t, pg in titles:
                if frag in t.replace("\n", " "):
                    return "%d장" % pg
            raise KeyError("제목에 없는 조각: " + frag)

        text = re.sub(r"«(\d부|부록)»", sub_part, text)
        return re.sub(r"«장:([^»]+)»", sub_slide, text)

    def walk(o):
        if isinstance(o, str):
            return one(o)
        if isinstance(o, list):
            return [walk(x) for x in o]
        if isinstance(o, tuple):
            return tuple(walk(x) for x in o)
        if isinstance(o, dict):
            return {k: walk(v) for k, v in o.items()}
        return o

    return walk(slides)


SLIDES = _resolve_refs(SLIDES)

# ── 색 ────────────────────────────────────────────────────────────────
BG      = "F6F8F9"   # 차가운 기운이 살짝 도는 흰색 — 인쇄해도 뜨지 않는다
INK     = "0E1A1F"   # 순수 검정 대신 짙은 청록빛 먹
MUTED   = "5C6E76"
ACCENT  = "0E8C7F"   # 앱의 청록을 인쇄에서도 읽히도록 낮춘 값
SOFT    = "E4F0EE"
WARN    = "A8530C"
LINE    = "D9E2E5"
PANEL   = "0B1418"   # 어두운 카드 (표지·강조 장)
PANEL_INK = "F2F7F8"

FONT = "Apple SD Gothic Neo"   # ⚠️ 윈도우에서 열면 '맑은 고딕'으로 바꿔 주세요


def _split_warn(t):
    """⚠️ 로 시작하는 조각을 따로 떼어 강조색을 준다."""
    return t.startswith("⚠️")


# ══════════════════════════════════════════════════════════════ HTML
def esc(t):
    return html.escape(t).replace("\n", "<br>")


def h_text(t):
    """⚠️ 와 `코드` 만 살려서 표시한다."""
    t = html.escape(t)
    t = re.sub(r"`([^`]+)`", r"<code>\1</code>", t)
    t = t.replace("⚠️", '<span class="w">⚠️</span>')
    return t.replace("\n", "<br>")


def h_slide(s):
    k = s["kind"]
    img = s.get("img")
    # ⚠️ 파일이 아직 없으면 깨진 그림 대신 "여기에 넣으라"는 판을 그린다.
    #    (브랜드 시트처럼 PD 가 나중에 넣는 그림이 있다)
    if img and not os.path.exists(os.path.join(IMG, img)):
        imgtag = ('<div class="missing">이미지 없음 — <code>docs/talk/img/%s</code> 로 '
                  '넣고 <code>python3 build.py</code></div>' % img)
    elif img:
        imgtag = '<img src="img/%s" alt="">' % img
    else:
        imgtag = ""

    if k == "cover":
        return f"""<section class="s cover">
  <div class="cov-bg">{imgtag}</div>
  <div class="cov-in">
    <img class="cov-logo" src="logo.png" alt="earthus">
    <h1>{h_text(s['title'])}</h1>
    <p class="sub">{h_text(s['sub'])}</p>
    <ul class="meta">{''.join('<li>%s</li>' % h_text(m) for m in s['meta'])}</ul>
  </div>
</section>"""

    if k == "section":
        return f"""<section class="s sect">
  <div class="sect-in">
    <span class="part">{h_text(s['part'])}</span>
    <h2>{h_text(s['title'])}</h2>
    {'<p class="sub">%s</p>' % h_text(s['sub']) if s.get('sub') else ''}
    {'<p class="note">%s</p>' % h_text(s['note']) if s.get('note') else ''}
  </div>
</section>"""

    if k == "big":
        return f"""<section class="s big">
  <div class="big-in">
    {'<span class="eyebrow">%s</span>' % h_text(s['eyebrow']) if s.get('eyebrow') else ''}
    <h2>{h_text(s['title'])}</h2>
    {'<p class="note">%s</p>' % h_text(s['note']) if s.get('note') else ''}
  </div>
</section>"""

    if k == "stats":
        cells = "".join(
            '<div class="stat"><b>%s</b><span>%s</span></div>' % (h_text(v), h_text(l))
            for v, l in s["stats"])
        return f"""<section class="s">
  <h3>{h_text(s['title'])}</h3>
  <div class="stats">{cells}</div>
  {'<p class="note">%s</p>' % h_text(s['note']) if s.get('note') else ''}
</section>"""

    if k == "bullets":
        lis = "".join('<li%s>%s</li>' % (' class="w-li"' if _split_warn(i) else '', h_text(i))
                      for i in s["items"])
        body = f'<ul class="bul">{lis}</ul>'
        if img:
            body = f'<div class="two"><div>{body}</div><figure class="shot">{imgtag}</figure></div>'
        return f"""<section class="s">
  <h3>{h_text(s['title'])}</h3>
  {body}
  {'<p class="note">%s</p>' % h_text(s['note']) if s.get('note') else ''}
</section>"""

    if k == "split":
        lis = "".join("<li>%s</li>" % h_text(i) for i in s["items"])
        return f"""<section class="s">
  <h3>{h_text(s['title'])}</h3>
  <div class="two">
    <div><ul class="bul">{lis}</ul></div>
    <figure class="shot">{imgtag}</figure>
  </div>
  {'<p class="note">%s</p>' % h_text(s['note']) if s.get('note') else ''}
</section>"""

    if k == "shot":
        return f"""<section class="s">
  <h3>{h_text(s['title'])}</h3>
  <figure class="shot big-shot">{imgtag}</figure>
  <p class="cap">{h_text(s['cap'])}</p>
</section>"""

    if k == "case":
        rows = [("이렇게 보였다", s["seen"]), ("실제로는", s["real"]), ("그래서", s["did"])]
        rh = "".join(
            f'<div class="row"><span class="lab">{lab}</span><p>{h_text(txt)}</p></div>'
            for lab, txt in rows)
        side = f'<figure class="shot side">{imgtag}</figure>' if img else ""
        return f"""<section class="s case">
  <span class="eyebrow">사례 {s['n']}</span>
  <h3>{h_text(s['title'])}</h3>
  <div class="case-body{' two' if img else ''}">
    <div class="rows">{rh}</div>
    {side}
  </div>
  <p class="tag">{h_text(s['tag'])}</p>
</section>"""

    if k == "case_lite":
        rh = "".join(f'<div class="row"><p class="lead">{h_text(a)}</p>'
                     f'<p class="sub2">{h_text(b)}</p></div>' for a, b in s["lines"])
        return f"""<section class="s">
  <h3>{h_text(s['title'])}</h3>
  <div class="rows lite">{rh}</div>
</section>"""

    if k == "models":
        cards = ""
        for name, when, items, why in s["rows"]:
            lis = "".join("<li>%s</li>" % h_text(i) for i in items)
            cards += (f'<div class="mcard"><b>{h_text(name)}</b>'
                      f'<span class="when">{h_text(when)}</span>'
                      f'<ul>{lis}</ul><p class="why">{h_text(why)}</p></div>')
        return f"""<section class="s">
  <h3>{h_text(s['title'])}</h3>
  <div class="models">{cards}</div>
  {'<p class="note">%s</p>' % h_text(s['note']) if s.get('note') else ''}
</section>"""

    return '<section class="s"><h3>?</h3></section>'


CSS = """
*,*::before,*::after{box-sizing:border-box}
html,body{margin:0;padding:0}
body{
  background:#%(bg)s; color:#%(ink)s;
  font-family:'Apple SD Gothic Neo','Pretendard','Noto Sans KR',
              -apple-system,BlinkMacSystemFont,'Malgun Gothic',sans-serif;
  -webkit-font-smoothing:antialiased;
  font-feature-settings:'tnum' 1;
}
code{font-family:ui-monospace,'SF Mono',Menlo,monospace;font-size:.9em;
     background:#%(soft)s;padding:.08em .34em;border-radius:5px}
.w{color:#%(warn)s;font-weight:700}

/* 한 장 = 1600×900 고정, 창 크기에는 축척으로 맞춘다.
   ⚠️ 글자만 줄이면 배치가 어긋난다 — 장을 통째로 줄여야 어느 창에서도 같은 그림이 나온다. */
:root{--z:1}
.s{
  position:relative; width:1600px; height:900px;
  margin:0 auto; padding:64px 76px 76px;
  background:#%(bg)s; display:none; flex-direction:column;
  border-radius:22px; overflow:hidden; zoom:var(--z);
}
.s.on{display:flex}
body:not(.all){display:flex;align-items:center;justify-content:center;
  min-height:100vh;overflow:hidden}
body.all{padding:20px 0}
body.all .s{display:flex;margin-bottom:22px;box-shadow:0 2px 14px rgba(14,26,31,.10)}

h1{font-size:58px;line-height:1.16;margin:0 0 18px;font-weight:800;letter-spacing:-.02em;
   text-wrap:balance}
h2{font-size:52px;line-height:1.2;margin:0;font-weight:800;letter-spacing:-.02em;
   text-wrap:balance}
h3{font-size:38px;line-height:1.25;margin:0 0 30px;font-weight:750;letter-spacing:-.018em;
   text-wrap:balance}
.eyebrow{display:inline-block;font-size:15px;font-weight:700;letter-spacing:.13em;
  text-transform:uppercase;color:#%(accent)s;margin-bottom:14px}
.sub{font-size:23px;color:#%(muted)s;margin:0 0 8px;line-height:1.5}
.note{font-size:19px;color:#%(muted)s;line-height:1.55;margin:26px 0 0;
      padding-top:18px;border-top:1px solid #%(line)s}
.cap{font-size:18px;color:#%(muted)s;line-height:1.55;margin:18px 0 0}

/* 표지 */
.cover{padding:0;background:#%(panel)s;color:#%(pink)s}
.cov-bg{position:absolute;inset:0}
.cov-bg img{width:100%%;height:100%%;object-fit:cover;opacity:.55}
.cov-bg::after{content:'';position:absolute;inset:0;background:linear-gradient(100deg,rgba(11,20,24,.92) 0%%,rgba(11,20,24,.72) 46%%,rgba(11,20,24,.30) 100%%)}
.cov-in{position:relative;margin:auto 0;padding:0 90px;max-width:1080px}
.cov-logo{height:40px;margin-bottom:34px;filter:brightness(0) invert(1);opacity:.95}
.cover h1{font-size:64px}
.cover .sub{color:#BFD6D8;font-size:24px}
.meta{list-style:none;display:flex;gap:26px;flex-wrap:wrap;padding:0;margin:34px 0 0;
  font-size:17px;color:#93AEB2}
.meta li:not(:last-child)::after{content:'';margin-left:26px;border-left:1px solid #35494F}

/* 장 표지 */
.sect{background:#%(panel)s;color:#%(pink)s}
.sect-in{margin:auto;max-width:1120px}
.sect h2{font-size:76px;color:#fff}
.part{display:block;font-size:19px;letter-spacing:.2em;color:#%(accent)s;
  font-weight:800;margin-bottom:16px}
.sect .sub{color:#BFD6D8;font-size:26px;margin-top:16px}
.sect .note{border-color:#283A40;color:#93AEB2}

/* 큰 문장 */
.big-in{margin:auto 0;max-width:1180px}
.big h2{font-size:66px}
.big{border-left:10px solid #%(accent)s}

/* ⚠️ 본문은 남는 세로 공간을 채우고 가운데로 온다.
   안 그러면 내용이 위쪽에만 몰리고 아래가 텅 빈다. */
.stats,.bul,.two,.rows,.models,.big-shot,.case-body{flex:1;min-height:0}
.case-body:not(.two){display:flex;align-items:center}
.case-body:not(.two)>.rows{width:100%%}
.stats,.models{align-content:center}
.bul,.rows{justify-content:center}
.two{align-items:center}
.note{margin-top:auto}

/* 숫자 */
.stats{display:grid;grid-template-columns:repeat(3,1fr);gap:20px;margin-top:6px}
.stat{background:#fff;border:1px solid #%(line)s;border-radius:16px;padding:26px 28px}
.stat b{display:block;font-size:46px;font-weight:800;letter-spacing:-.03em;
  color:#%(accent)s;line-height:1.05;font-variant-numeric:tabular-nums}
.stat span{display:block;margin-top:8px;font-size:17px;color:#%(muted)s;line-height:1.4}

/* 목록 */
.bul{list-style:none;padding:0;margin:0;display:flex;flex-direction:column;gap:18px}
.bul li{position:relative;padding-left:32px;font-size:26px;line-height:1.55}
.bul li::before{content:'';position:absolute;left:2px;top:.62em;width:11px;height:2px;
  background:#%(accent)s;border-radius:2px}
.bul li.w-li::before{background:#%(warn)s}

/* 두 칸 */
.two{display:grid;grid-template-columns:1fr 1fr;gap:44px;align-items:start;flex:1;min-height:0}
.shot{margin:0;border-radius:14px;overflow:hidden;background:#%(panel)s;
  border:1px solid #%(line)s;box-shadow:0 6px 22px rgba(14,26,31,.10)}
.shot img{display:block;width:100%%;height:auto}
.missing{display:flex;align-items:center;justify-content:center;min-height:340px;padding:40px;color:#93AEB2;font-size:20px;text-align:center;line-height:1.7}
.missing code{background:rgba(255,255,255,.08);color:#BFD6D8}
.big-shot{flex:1;min-height:0;display:flex}
.big-shot img{margin:auto;max-height:100%%;width:auto;max-width:100%%;object-fit:contain}

/* 사례 */
.case .eyebrow{margin-bottom:10px}
.rows{display:flex;flex-direction:column;gap:26px}
.row{display:grid;grid-template-columns:150px 1fr;gap:22px;align-items:start}
.lab{font-size:17px;font-weight:800;color:#%(accent)s;padding-top:5px;letter-spacing:-.01em}
.row p{margin:0;font-size:25px;line-height:1.55}
.rows.lite .row{display:block;padding-left:26px;border-left:3px solid #%(accent)s}
.lead{font-size:30px!important;font-weight:750;line-height:1.35!important}
.sub2{margin-top:6px!important;color:#%(muted)s;font-size:20px!important}
.tag{align-self:flex-start;margin:26px 0 0;padding:14px 22px;background:#%(soft)s;border-radius:12px;
  font-size:19px;font-weight:650;color:#%(ink)s;display:inline-block}

/* 모델 */
.models{display:grid;grid-template-columns:1fr 1fr;gap:26px}
.mcard{background:#fff;border:1px solid #%(line)s;border-radius:18px;padding:30px 32px}
.mcard b{font-size:34px;font-weight:800;color:#%(accent)s;letter-spacing:-.02em}
.when{display:block;margin:6px 0 20px;font-size:19px;color:#%(muted)s;font-weight:600}
.mcard ul{list-style:none;padding:0;margin:0;display:flex;flex-direction:column;gap:11px}
.mcard li{font-size:22px;padding-left:22px;position:relative;line-height:1.45}
.mcard li::before{content:'';position:absolute;left:0;top:.6em;width:9px;height:2px;
  background:#%(accent)s;border-radius:2px}
.why{margin:22px 0 0;padding-top:16px;border-top:1px solid #%(line)s;
  font-size:18px;color:#%(muted)s;line-height:1.5}

/* 아래 표시 */
.pager{position:fixed;right:22px;bottom:18px;font-size:14px;color:#%(muted)s;
  background:rgba(246,248,249,.9);padding:7px 13px;border-radius:999px;
  border:1px solid #%(line)s;font-variant-numeric:tabular-nums}
.hint{position:fixed;left:22px;bottom:18px;font-size:13px;color:#%(muted)s;background:rgba(246,248,249,.9);padding:7px 13px;border-radius:999px;border:1px solid #%(line)s}
body.all .pager,body.all .hint{display:none}
.s::after{content:attr(data-n);position:absolute;right:34px;bottom:22px;
  font-size:14px;color:#%(muted)s;font-variant-numeric:tabular-nums}
.cover::after,.sect::after{color:#6E888D}

@media print{
  @page{size:1600px 900px;margin:0}
  body{background:#fff}
  .pager,.hint{display:none}
  body:not(.all){display:block;overflow:visible}
  .s{display:flex!important;width:1600px;height:900px;margin:0;border-radius:0;
     break-after:page;box-shadow:none;zoom:1}
  .stat,.mcard{background:#fff}
}
""" % dict(bg=BG, ink=INK, muted=MUTED, accent=ACCENT, soft=SOFT, warn=WARN,
           line=LINE, panel=PANEL, pink=PANEL_INK)

JS = """
const S=[...document.querySelectorAll('.s')];let i=0;
function fit(){
  const all=document.body.classList.contains('all');
  const k=all ? Math.min(1,(innerWidth-40)/1600)
              : Math.min(innerWidth/1600, innerHeight/900);
  document.documentElement.style.setProperty('--z', k.toFixed(4));
}
addEventListener('resize',fit);
const pg=document.querySelector('.pager');
function go(n){i=Math.max(0,Math.min(S.length-1,n));
  S.forEach((s,k)=>s.classList.toggle('on',k===i));
  pg.textContent=(i+1)+' / '+S.length;location.hash=i+1;}
addEventListener('keydown',e=>{
  if(document.body.classList.contains('all'))return;
  if(e.key==='ArrowRight'||e.key===' '||e.key==='PageDown'){e.preventDefault();go(i+1)}
  if(e.key==='ArrowLeft'||e.key==='PageUp'){e.preventDefault();go(i-1)}
  if(e.key==='Home')go(0); if(e.key==='End')go(S.length-1);
  if(e.key==='a'||e.key==='A'){document.body.classList.toggle('all');fit();
    if(!document.body.classList.contains('all'))go(i);}
});
addEventListener('click',e=>{if(e.target.closest('a'))return;
  if(!document.body.classList.contains('all'))go(i+1)});
if(location.search.indexOf('all')>=0)document.body.classList.add('all');
fit();
go(location.hash?parseInt(location.hash.slice(1))-1||0:0);
"""


def build_html(path):
    body = "\n".join(h_slide(s) for s in SLIDES)
    # 장 번호를 넣어 준다
    n = [0]

    def num(m):
        n[0] += 1
        return '<section class="s%s" data-n="%d"' % (m.group(1), n[0])

    body = re.sub(r'<section class="s([^"]*)"', num, body)

    doc = f"""<!doctype html>
<meta charset="utf-8">
<title>earthus — 지구를 믿을 수 있게 만드는 법</title>
<meta name="viewport" content="width=device-width,initial-scale=1">
<style>{CSS}</style>
{body}
<div class="pager"></div>
<div class="hint">← → 넘기기 · A 전체보기 · ⌘P 인쇄</div>
<script>{JS}</script>
"""
    io.open(path, "w", encoding="utf-8").write(doc)
    return len(SLIDES)


# ══════════════════════════════════════════════════════════════ PPTX
def build_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches as In, Pt
    from pptx.dml.color import RGBColor as C
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    def rgb(h):
        return C(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    prs = Presentation()
    prs.slide_width, prs.slide_height = In(13.333), In(7.5)
    W, H = 13.333, 7.5
    M = 0.85                      # 좌우 여백
    CW = W - M * 2                # 본문 폭
    blank = prs.slide_layouts[6]

    def bg(sl, hexv):
        f = sl.background.fill
        f.solid()
        f.fore_color.rgb = rgb(hexv)

    def box(sl, x, y, w, h, mid=False):
        tb = sl.shapes.add_textbox(In(x), In(y), In(w), In(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        if mid:
            # ⚠️ 가운데 정렬을 안 하면 내용이 위에만 몰리고 아래가 빈다
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        return tf

    def para(tf, text, size, color, bold=False, space_before=0, space_after=0,
             line=1.35, first=False):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.line_spacing = line
        p.space_before = Pt(space_before)
        p.space_after = Pt(space_after)
        # ⚠️ 조각만 강조색으로 나눠 넣는다
        parts = re.split(r"(⚠️)", text)
        emph = False
        for seg in parts:
            if seg == "":
                continue
            r = p.add_run()
            r.text = seg
            r.font.size = Pt(size)
            r.font.name = FONT
            r.font.bold = bold or (seg == "⚠️")
            r.font.color.rgb = rgb(WARN if (seg == "⚠️" or emph) else color)
            if seg == "⚠️":
                emph = True
        return p

    def rect(sl, x, y, w, h, fill, line=None, radius=True):
        from pptx.enum.shapes import MSO_SHAPE
        shp = sl.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
            In(x), In(y), In(w), In(h))
        shp.shadow.inherit = False
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        if line:
            shp.line.color.rgb = rgb(line)
            shp.line.width = Pt(0.75)
        else:
            shp.line.fill.background()
        if radius:
            try:
                shp.adjustments[0] = 0.06
            except Exception:
                pass
        shp.text_frame.word_wrap = True
        return shp

    def pic(sl, name, x, y, w):
        p = os.path.join(IMG, name)
        if not os.path.exists(p):
            return None
        return sl.shapes.add_picture(p, In(x), In(y), width=In(w))

    def title(sl, text, y=1.05, size=32, color=INK):
        tf = box(sl, M, y, CW, 1.4)
        para(tf, text, size, color, bold=True, line=1.2, first=True)
        return tf

    def eyebrow(sl, text, y=0.62):
        tf = box(sl, M, y, CW, 0.32)
        para(tf, text, 12.5, ACCENT, bold=True, first=True)

    def note(sl, text, y=None):
        yy = y if y is not None else H - 1.28
        ln = sl.shapes.add_shape(1, In(M), In(yy - 0.16), In(CW), In(0.012))
        ln.fill.solid(); ln.fill.fore_color.rgb = rgb(LINE)
        ln.line.fill.background(); ln.shadow.inherit = False
        tf = box(sl, M, yy, CW, 1.0)
        para(tf, text, 13, MUTED, line=1.4, first=True)

    def pagenum(sl, n, color=MUTED):
        tf = box(sl, W - M - 1.2, H - 0.58, 1.2, 0.3)
        p = para(tf, str(n), 10.5, color, first=True)
        p.alignment = PP_ALIGN.RIGHT

    # ── 장별 ───────────────────────────────────────────────────
    for n, s in enumerate(SLIDES, 1):
        sl = prs.slides.add_slide(blank)
        k = s["kind"]
        img = s.get("img")

        if k == "cover":
            bg(sl, PANEL)
            if img:
                # 배경 사진을 가득 깔고 그 위에 어두운 판을 덮는다
                p = pic(sl, img, 0, 0, W)
                if p:
                    p.top = In((H - p.height.inches) / 2)
                sh = rect(sl, -0.1, -0.1, W + 0.2, H + 0.2, PANEL, radius=False)
                sh.fill.fore_color.rgb = rgb(PANEL)
                sh.fill.transparency = 0.42
            if os.path.exists(LOGO):
                sl.shapes.add_picture(LOGO, In(M + 0.15), In(1.5), height=In(0.38))
            tf = box(sl, M + 0.15, 2.15, CW - 1.6, 2.0)
            para(tf, s["title"], 40, "FFFFFF", bold=True, line=1.14, first=True)
            tf2 = box(sl, M + 0.15, 4.15, CW - 1.6, 0.9)
            para(tf2, s["sub"], 15.5, "BFD6D8", line=1.5, first=True)
            tf3 = box(sl, M + 0.15, 5.25, CW - 1.6, 1.2)
            for i, m in enumerate(s["meta"]):
                para(tf3, m, 12, "93AEB2", first=(i == 0), space_before=3)
            pagenum(sl, n, "6E888D")
            continue

        if k == "section":
            bg(sl, PANEL)
            tf = box(sl, M + 0.15, 2.5, CW - 1.0, 0.4)
            para(tf, s["part"], 13, ACCENT, bold=True, first=True)
            tf2 = box(sl, M + 0.15, 3.0, CW - 1.0, 1.3)
            para(tf2, s["title"], 50, "FFFFFF", bold=True, line=1.12, first=True)
            yy = 4.45
            if s.get("sub"):
                tf3 = box(sl, M + 0.15, yy, CW - 1.6, 0.7)
                para(tf3, s["sub"], 17, "BFD6D8", line=1.4, first=True)
                yy += 0.75
            if s.get("note"):
                tf4 = box(sl, M + 0.15, yy, CW - 1.6, 0.8)
                para(tf4, s["note"], 12.5, "93AEB2", line=1.45, first=True)
            pagenum(sl, n, "6E888D")
            continue

        bg(sl, BG)
        pagenum(sl, n)

        if k == "big":
            bar = rect(sl, 0, 0, 0.09, H, ACCENT, radius=False)
            bar.line.fill.background()
            if s.get("eyebrow"):
                eyebrow(sl, s["eyebrow"], 2.15)
            tf = box(sl, M, 2.6, CW - 1.2, 2.4)
            para(tf, s["title"], 40, INK, bold=True, line=1.18, first=True)
            if s.get("note"):
                note(sl, s["note"], H - 1.35)
            continue

        if k == "stats":
            title(sl, s["title"], 0.72, 30)
            cols, gap = 3, 0.24
            cw = (CW - gap * (cols - 1)) / cols
            ch = 1.55
            nrow = (len(s["stats"]) + cols - 1) // cols
            top, bot = 1.95, (H - 1.5 if s.get("note") else H - 0.7)
            y0 = top + max(0, (bot - top - (nrow * ch + (nrow - 1) * gap)) / 2)
            for idx, (v, l) in enumerate(s["stats"]):
                r, c = divmod(idx, cols)
                x = M + c * (cw + gap)
                y = y0 + r * (ch + gap)
                rect(sl, x, y, cw, ch, "FFFFFF", LINE)
                tf = box(sl, x + 0.3, y + 0.26, cw - 0.6, 0.62)
                para(tf, v, 31, ACCENT, bold=True, line=1.0, first=True)
                tf2 = box(sl, x + 0.3, y + 0.94, cw - 0.6, 0.5)
                para(tf2, l, 11.5, MUTED, line=1.3, first=True)
            if s.get("note"):
                note(sl, s["note"])
            continue

        if k in ("bullets", "split"):
            title(sl, s["title"], 0.72, 30)
            has_note = bool(s.get("note"))
            bottom = (H - 1.55) if has_note else (H - 0.75)
            top = 1.85
            if img:
                bw = CW * 0.46
                pw = CW - bw - 0.5
                ph = pw * 9 / 16
                pic(sl, img, M + bw + 0.5, top + max(0, (bottom - top - ph) / 2), pw)
                tf = box(sl, M, top, bw, bottom - top, mid=True)
            else:
                tf = box(sl, M, top, CW - 0.6, bottom - top, mid=True)
            size = 16 if img else 19
            for i, it in enumerate(s["items"]):
                para(tf, "·  " + it, size, INK, line=1.45,
                     space_before=0 if i == 0 else 13, first=(i == 0))
            if has_note:
                note(sl, s["note"])
            continue

        if k == "shot":
            title(sl, s["title"], 0.72, 30)
            # ⚠️ 그림 높이 + 설명 글이 장 아래를 넘지 않게 폭부터 잡는다
            ptop, cap_h, gap = 1.85, 0.82, 0.2
            pw = min(CW, (H - ptop - cap_h - gap) * 16 / 9)
            px = (W - pw) / 2
            if pic(sl, img, px, ptop, pw) is None:
                ph = pw * 9 / 16
                rect(sl, px, ptop, pw, ph, "FFFFFF", LINE)
                tfm = box(sl, px + 0.5, ptop, pw - 1.0, ph, mid=True)
                para(tfm, "이미지 없음 — docs/talk/img/%s 로 넣고 python3 build.py" % img,
                     15, MUTED, line=1.5, first=True)
            cy = ptop + pw * 9 / 16 + gap
            tf = box(sl, px, cy, pw, cap_h)
            para(tf, s["cap"], 12.5, MUTED, line=1.45, first=True)
            continue

        if k == "case":
            eyebrow(sl, "사례 %d" % s["n"], 0.62)
            title(sl, s["title"], 1.0, 30)
            rows = [("이렇게 보였다", s["seen"]), ("실제로는", s["real"]), ("그래서", s["did"])]
            top, bot = 1.95, H - 1.35
            wide = CW if not img else CW * 0.53
            tf = box(sl, M, top, wide - 0.2, bot - top, mid=True)
            first = True
            for lab, txt in rows:
                para(tf, lab, 11.5, ACCENT, bold=True, line=1.1,
                     space_before=0 if first else 15, space_after=3, first=first)
                para(tf, txt, 14 if img else 15.5, INK, line=1.42)
                first = False
            if img:
                pw = CW - wide - 0.25
                pic(sl, img, M + wide + 0.25, top + max(0, (bot - top - pw * 9 / 16) / 2), pw)
            tag = rect(sl, M, H - 1.12, min(CW, len(s["tag"]) * 0.148 + 0.7), 0.62, SOFT)
            tf = box(sl, M + 0.28, H - 0.97, CW - 0.6, 0.5)
            para(tf, s["tag"], 13, INK, bold=True, first=True)
            continue

        if k == "case_lite":
            title(sl, s["title"], 0.72, 30)
            y = 2.0 + max(0, ((H - 0.8) - 2.0 - len(s["lines"]) * 1.42) / 2)
            for a, b in s["lines"]:
                bar = rect(sl, M, y, 0.045, 1.05, ACCENT, radius=False)
                bar.line.fill.background()
                tf = box(sl, M + 0.32, y, CW - 0.5, 0.55)
                para(tf, a, 22, INK, bold=True, line=1.25, first=True)
                tf2 = box(sl, M + 0.32, y + 0.55, CW - 0.5, 0.5)
                para(tf2, b, 14, MUTED, line=1.35, first=True)
                y += 1.42
            continue

        if k == "models":
            title(sl, s["title"], 0.72, 30)
            gap = 0.34
            cw = (CW - gap) / 2
            ctop = 1.9
            ch = (H - 1.5 if s.get("note") else H - 0.7) - ctop
            for idx, (name, when, items, why) in enumerate(s["rows"]):
                x = M + idx * (cw + gap)
                rect(sl, x, ctop, cw, ch, "FFFFFF", LINE)
                tf = box(sl, x + 0.34, ctop + 0.3, cw - 0.68, 0.55)
                para(tf, name, 26, ACCENT, bold=True, line=1.05, first=True)
                tf2 = box(sl, x + 0.34, ctop + 0.92, cw - 0.68, 0.35)
                para(tf2, when, 13.5, MUTED, bold=True, first=True)
                tf3 = box(sl, x + 0.34, ctop + 1.45, cw - 0.68, ch - 2.35, mid=True)
                for i, it in enumerate(items):
                    para(tf3, "·  " + it, 14.5, INK, line=1.35,
                         space_before=0 if i == 0 else 10, first=(i == 0))
                tf4 = box(sl, x + 0.34, ctop + ch - 0.78, cw - 0.68, 0.6)
                para(tf4, why, 12, MUTED, line=1.4, first=True)
            if s.get("note"):
                note(sl, s["note"])
            continue

    prs.save(path)
    return len(prs.slides.__iter__.__self__._sldIdLst)


if __name__ == "__main__":
    # 표지 로고를 발표 폴더로 복사해 둔다 (HTML 이 상대경로로 쓴다)
    dst = os.path.join(HERE, "logo.png")
    if os.path.exists(LOGO) and not os.path.exists(dst):
        io.open(dst, "wb").write(io.open(LOGO, "rb").read())

    n1 = build_html(os.path.join(HERE, "index.html"))
    build_pptx(os.path.join(HERE, "earthus-talk.pptx"))
    print("index.html          %d장" % n1)
    print("earthus-talk.pptx   %d장" % n1)
